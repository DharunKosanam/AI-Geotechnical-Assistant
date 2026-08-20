"""
File processing service: image conversion + unified multi-format text extraction.

Extractors all return List[Tuple[int, str, bool]] where:
    page_number   — 1-indexed logical page (sheet, slide, or PDF page)
    text          — extracted text, with embedded markdown headers so the v2
                    chunker can pick up section structure
    ocr_extracted — True when this page's text came from OCR (image input,
                    OCR'd PDF page, or OCR'd embedded PDF image)

All optional dependencies are soft-imported; calling an extractor for a
format whose library isn't installed raises a clear RuntimeError instead of
failing at import time.
"""
import csv
import datetime as _dt
import io
import os
import shutil
from functools import lru_cache
from itertools import zip_longest
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Sequence, Tuple

from PIL import Image
import img2pdf
from fastapi import HTTPException, status

from app.core.config import (
    CHUNK_TARGET_SIZE,
    OCR_ENABLED,
    OCR_MIN_TEXT_LEN,
    TESSERACT_PATH,
    XLSX_MAX_ROWS_PER_SHEET,
)


# ---------------------------------------------------------------------------
# Optional dependency probes — each format degrades gracefully if missing
# ---------------------------------------------------------------------------
try:
    import pytesseract  # type: ignore
    _PYTESSERACT_IMPORTED = True
    if TESSERACT_PATH:
        pytesseract.pytesseract.tesseract_cmd = TESSERACT_PATH
except Exception:
    pytesseract = None  # type: ignore
    _PYTESSERACT_IMPORTED = False

try:
    from docx import Document as _DocxDocument  # type: ignore
    _DOCX_AVAILABLE = True
except Exception:
    _DocxDocument = None  # type: ignore
    _DOCX_AVAILABLE = False

try:
    import openpyxl  # type: ignore
    _OPENPYXL_AVAILABLE = True
except Exception:
    openpyxl = None  # type: ignore
    _OPENPYXL_AVAILABLE = False

try:
    import pandas as _pd  # type: ignore
    _PANDAS_AVAILABLE = True
except Exception:
    _pd = None  # type: ignore
    _PANDAS_AVAILABLE = False

try:
    from pptx import Presentation as _Presentation  # type: ignore
    _PPTX_AVAILABLE = True
except Exception:
    _Presentation = None  # type: ignore
    _PPTX_AVAILABLE = False


# Public list of file extensions handled by the unified extractor
SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx",
    ".xlsx", ".xls", ".csv",
    ".png", ".jpg", ".jpeg", ".tiff", ".tif",
    ".pptx",
}

_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".tiff", ".tif"}

# Image formats the VISION path accepts for direct upload (gated on
# VISION_EXTRACTION_ENABLED at the call sites -- this set is just the format
# registry). Deliberately excludes TIFF: it keeps its legacy OCR-only path.
# .webp is NOT in SUPPORTED_EXTENSIONS: with vision off it stays rejected
# exactly as today, so flag-off behavior is byte-identical.
VISION_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp"}

# The formats that carry their own text and so never depend on OCR. Used in
# user-facing messages when we have to turn an image or a scan away.
TEXT_FORMATS_LABEL = "PDF, DOCX, XLSX, XLS, CSV, PPTX"


class UnreadableDocumentError(Exception):
    """A file we could open and parse, but from which no text can be read, for
    a reason we can state plainly to the user (a scanned PDF with no text layer,
    an image on a server without OCR, ...).

    Distinct from a generic failure so the upload chip shows what is actually
    wrong and what to do about it, instead of "no extractable text" or, worse,
    "Upload failed (500)". Its message is USER-FACING -- keep it short, name the
    file, and say what to do next.
    """


def is_image_file(filename: str) -> bool:
    return get_file_type(filename) in _IMAGE_EXTS


def get_file_type(filename: str) -> str:
    """Return the lowercase extension including dot (e.g. '.pdf')."""
    return Path(filename).suffix.lower()


def is_supported_file(filename: str) -> bool:
    return get_file_type(filename) in SUPPORTED_EXTENSIONS


@lru_cache(maxsize=1)
def _tesseract_binary_present() -> bool:
    """True only when the tesseract EXECUTABLE can actually be run.

    The pytesseract package importing says nothing about the binary: pytesseract
    is a thin wrapper that shells out to `tesseract`, and importing it succeeds
    on a machine that has never had OCR installed. Trusting the import meant
    every sparse PDF page and every embedded image entered the OCR path and
    raised TesseractNotFoundError inside a swallow-and-return-"" handler -- pure
    wasted work (measured: ~28 s of ingest CPU on a 120-page scan) that ended in
    a generic "no extractable text" failure instead of a clear one.

    Probed once per process: the answer only changes when tesseract is installed
    or removed, which needs a backend restart anyway.
    """
    if not _PYTESSERACT_IMPORTED:
        return False
    cmd = TESSERACT_PATH or getattr(pytesseract.pytesseract, "tesseract_cmd", "tesseract")
    if shutil.which(cmd) is None and not os.path.isfile(cmd):
        return False
    try:
        # Definitive: actually invoke it. Catches an unusable/broken install
        # that a which() hit alone would happily report as present.
        pytesseract.get_tesseract_version()
        return True
    except Exception as e:
        print(f"[OCR] tesseract binary found at {cmd!r} but not usable: {e}")
        return False


def tesseract_available() -> bool:
    return _tesseract_binary_present() and OCR_ENABLED


# ---------------------------------------------------------------------------
# OCR helper — used by image extractor AND by PDF OCR fallback in rag_service
# ---------------------------------------------------------------------------
def ocr_image_bytes(image_bytes: bytes) -> str:
    """OCR an image given as raw bytes. Returns empty string on any failure."""
    if not tesseract_available():
        return ""
    try:
        img = Image.open(io.BytesIO(image_bytes))
        return pytesseract.image_to_string(img) or ""
    except Exception as e:
        print(f"[OCR ERROR] {e}")
        return ""


def ocr_pil_image(pil_image) -> str:
    """OCR an already-open PIL.Image. Returns empty string on failure."""
    if not tesseract_available():
        return ""
    try:
        return pytesseract.image_to_string(pil_image) or ""
    except Exception as e:
        print(f"[OCR ERROR] {e}")
        return ""


# ---------------------------------------------------------------------------
# Format-specific extractors
# ---------------------------------------------------------------------------
def extract_pages_from_docx(content: bytes) -> List[Tuple[int, str, bool]]:
    """Word doc — preserves heading styles as markdown headers; tables included."""
    if not _DOCX_AVAILABLE:
        raise RuntimeError("python-docx is not installed (pip install python-docx)")

    doc = _DocxDocument(io.BytesIO(content))
    parts: List[str] = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style_name = (para.style.name if para.style else "") or ""
        if style_name.startswith("Heading"):
            # Map 'Heading 2' -> '##'; default to '##' if level unparseable
            tail = style_name.split()[-1]
            level = int(tail) if tail.isdigit() else 2
            hashes = "#" * min(6, max(1, level))
            parts.append(f"\n{hashes} {text}\n")
        else:
            parts.append(text + "\n")

    for ti, table in enumerate(doc.tables, 1):
        parts.append(f"\n## Table {ti}\n")
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            parts.append(" | ".join(cells) + "\n")

    full = "".join(parts).strip()
    return [(1, full, False)] if full else []


# ---------------------------------------------------------------------------
# Spreadsheet extraction (v3-xlsx) — row-structured, shared by the live upload
# path AND the KB registry (kb_formats). Each spreadsheet row becomes ONE line
# with its columns in order, so row associations ("who booked what") survive
# chunking and retrieval.
# ---------------------------------------------------------------------------
def _format_cell(value: Any) -> str:
    """One cell as the readable value a person sees in Excel: dates as
    YYYY-MM-DD (time appended only when non-midnight), int-valued floats
    without the trailing .0. Newlines/pipes inside a cell are replaced so a
    row always stays one pipe-delimited line."""
    if value is None:
        return ""
    if isinstance(value, bool):
        return "TRUE" if value else "FALSE"
    if isinstance(value, _dt.datetime):
        if (value.hour, value.minute, value.second, value.microsecond) == (0, 0, 0, 0):
            return value.strftime("%Y-%m-%d")
        return value.strftime("%Y-%m-%d %H:%M")
    if isinstance(value, _dt.date):
        return value.strftime("%Y-%m-%d")
    if isinstance(value, _dt.time):
        return value.strftime("%H:%M") if value.second == 0 else value.strftime("%H:%M:%S")
    if isinstance(value, float):
        if value != value:  # NaN (pandas .xls path)
            return ""
        if value.is_integer() and abs(value) < 1e15:
            return str(int(value))
        return str(value)
    s = str(value)
    if "\n" in s:
        s = "; ".join(part.strip() for part in s.splitlines() if part.strip())
    return s.replace("|", "/").strip()


def _render_sheet_rows(
    sheet_name: str,
    raw_rows: Iterable[Sequence[Any]],
    *,
    kind: str = "Sheet",
    max_rows: int = None,  # type: ignore[assignment]
    block_budget: int = None,  # type: ignore[assignment]
) -> Tuple[Optional[str], Dict[str, Any]]:
    """Render one sheet's rows as row-per-line blocks.

    The first non-empty row is treated as the header. Rows are grouped into
    blocks of at most ``block_budget`` characters, and EVERY block restates the
    sheet name and the header line, so a chunk taken from the middle of a large
    table is still interpretable on its own. Blocks stay under the chunker's
    target size, so the v2 chunker never has to split inside one.

    Fully empty rows are skipped; fully empty TRAILING columns are trimmed
    (openpyxl often reports phantom width); blank cells inside the used width
    are kept as empty fields so column position carries meaning. Values are
    never deduplicated.

    Returns (text or None if the sheet is empty, meta) where meta carries
    ``header`` (list of header cells), ``rows`` (data rows rendered) and
    ``truncated_rows`` (data rows beyond the cap, 0 if none).
    """
    if max_rows is None:
        max_rows = XLSX_MAX_ROWS_PER_SHEET
    if block_budget is None:
        block_budget = CHUNK_TARGET_SIZE

    formatted: List[List[str]] = []
    truncated = 0
    width = 0
    for raw in raw_rows:
        cells = [_format_cell(c) for c in (raw or ())]
        # Trim this row's trailing blanks; a fully blank row is skipped.
        while cells and not cells[-1]:
            cells.pop()
        if not cells:
            continue
        # First non-empty row is the header and does not count against the cap.
        if len(formatted) - 1 >= max_rows:
            truncated += 1
            continue
        formatted.append(cells)
        width = max(width, len(cells))
    if not formatted:
        return None, {"header": [], "rows": 0, "truncated_rows": 0}

    header = formatted[0]
    data_rows = formatted[1:]

    def line(cells: List[str]) -> str:
        return "| " + " | ".join(cells + [""] * (width - len(cells))) + " |"

    header_line = line(header)
    blocks: List[str] = []
    block: List[str] = [f"## {kind}: {sheet_name}", header_line]
    block_len = sum(len(x) + 1 for x in block)
    for cells in data_rows:
        row_line = line(cells)
        if block_len + len(row_line) > block_budget and len(block) > 2:
            blocks.append("\n".join(block))
            block = [f"## {kind}: {sheet_name} (continued)", header_line]
            block_len = sum(len(x) + 1 for x in block)
        block.append(row_line)
        block_len += len(row_line) + 1
    blocks.append("\n".join(block))

    text = "\n\n".join(blocks)
    if truncated:
        text += (f"\n[Truncated: {truncated} additional non-empty row(s) beyond the "
                 f"{max_rows}-row limit (XLSX_MAX_ROWS_PER_SHEET) are not included]")
    return text, {"header": header, "rows": len(data_rows), "truncated_rows": truncated}


def _merge_formula_row(values_row: Sequence[Any], formulas_row: Sequence[Any]) -> Tuple[List[Any], int]:
    """Prefer the cached value; where a formula cell has NO cached value
    (data_only=True gave None), fall back to the formula string so the cell is
    not silently blank. Returns (merged row, fallback count)."""
    merged: List[Any] = []
    fallbacks = 0
    for v, f in zip_longest(values_row or (), formulas_row or ()):
        if v is None and f is not None:
            if isinstance(f, str) and f.startswith("="):
                merged.append(f)
                fallbacks += 1
                continue
            text = getattr(f, "text", None)  # openpyxl ArrayFormula
            if text is not None:
                merged.append(str(text))
                fallbacks += 1
                continue
        merged.append(v)
    return merged, fallbacks


def extract_xlsx_sheets(content: bytes, filename: str = "") -> Tuple[List[Tuple[int, str, bool]], Dict[str, Any]]:
    """Excel .xlsx — one logical page per sheet, one line per row (v3-xlsx).

    The workbook is read twice: data_only=True for cached formula values, and
    data_only=False so a formula cell whose cached value is missing (typical of
    files saved by openpyxl or by tools that skip recalculation) falls back to
    its formula string instead of vanishing.

    Returns (pages, meta); meta carries per-sheet headers plus truncation and
    formula-fallback counts for the KB registry and for logging.
    """
    if not _OPENPYXL_AVAILABLE:
        raise RuntimeError("openpyxl is not installed (pip install openpyxl)")

    wb_values = openpyxl.load_workbook(io.BytesIO(content), data_only=True, read_only=True)
    try:
        wb_formulas = openpyxl.load_workbook(io.BytesIO(content), data_only=False, read_only=True)
    except Exception:
        wb_formulas = None  # values-only degradation; never fail the upload for this

    pages: List[Tuple[int, str, bool]] = []
    meta: Dict[str, Any] = {"sheets": [], "headers": {}, "truncated_sheets": [],
                            "formula_fallback_cells": 0}
    try:
        for idx, sheet_name in enumerate(wb_values.sheetnames, 1):
            ws_v = wb_values[sheet_name]
            ws_f = wb_formulas[sheet_name] if wb_formulas is not None else None

            def rows():
                f_iter = ws_f.iter_rows(values_only=True) if ws_f is not None else iter(())
                for v_row, f_row in zip_longest(ws_v.iter_rows(values_only=True), f_iter):
                    if ws_f is None or f_row is None:
                        yield v_row or ()
                        continue
                    merged, n = _merge_formula_row(v_row or (), f_row)
                    meta["formula_fallback_cells"] += n
                    yield merged

            text, sheet_meta = _render_sheet_rows(sheet_name, rows())
            if text is None:
                continue
            meta["sheets"].append(sheet_name)
            meta["headers"][sheet_name] = sheet_meta["header"]
            if sheet_meta["truncated_rows"]:
                meta["truncated_sheets"].append(sheet_name)
            pages.append((idx, text, False))
    finally:
        wb_values.close()
        if wb_formulas is not None:
            wb_formulas.close()

    if meta["formula_fallback_cells"]:
        print(f"[XLSX] {filename or 'workbook'}: {meta['formula_fallback_cells']} formula "
              f"cell(s) had no cached value; indexed the formula text instead")
    return pages, meta


def extract_pages_from_xlsx(content: bytes, filename: str = "") -> List[Tuple[int, str, bool]]:
    """Excel .xlsx — one logical page per sheet (see extract_xlsx_sheets)."""
    pages, _meta = extract_xlsx_sheets(content, filename)
    return pages


def _decode_text_bytes(content: bytes) -> str:
    """UTF-8 (with BOM) first, latin-1 fallback: never raises on any bytes."""
    for enc in ("utf-8-sig", "utf-8"):
        try:
            return content.decode(enc)
        except UnicodeDecodeError:
            continue
    return content.decode("latin-1", errors="replace")


def extract_csv_rows(content: bytes, filename: str = "") -> Tuple[List[Tuple[int, str, bool]], Dict[str, Any]]:
    """CSV — one line per row, first row as the repeated header (v3-xlsx).

    stdlib csv handles ragged rows natively; cells are kept verbatim (a CSV is
    already text, so there are no serials or datetime reprs to normalise)."""
    name = os.path.basename(filename) or "CSV Data"
    reader = csv.reader(io.StringIO(_decode_text_bytes(content)))
    text, sheet_meta = _render_sheet_rows(name, reader, kind="CSV")
    meta = {"columns": sheet_meta["header"], "rows": sheet_meta["rows"],
            "truncated_rows": sheet_meta["truncated_rows"]}
    pages: List[Tuple[int, str, bool]] = [(1, text, False)] if text else []
    return pages, meta


def extract_pages_from_csv(content: bytes, filename: str = "") -> List[Tuple[int, str, bool]]:
    pages, _meta = extract_csv_rows(content, filename)
    return pages


def extract_pages_from_xls(content: bytes, filename: str = "") -> List[Tuple[int, str, bool]]:
    """Legacy .xls via pandas/xlrd (cached values only — xlrd has no formula
    text), rendered through the same row-per-line path as .xlsx."""
    if not _PANDAS_AVAILABLE:
        raise RuntimeError("pandas is not installed (pip install pandas xlrd)")
    sheets = _pd.read_excel(io.BytesIO(content), sheet_name=None, header=None)
    pages: List[Tuple[int, str, bool]] = []
    for idx, (name, df) in enumerate(sheets.items(), 1):
        raw_rows = df.astype(object).where(df.notna(), None).values.tolist()
        text, _meta = _render_sheet_rows(str(name), raw_rows)
        if text:
            pages.append((idx, text, False))
    return pages


def extract_pages_from_image(content: bytes, filename: str = "This image") -> List[Tuple[int, str, bool]]:
    """OCR a standalone image file. Marked as OCR-extracted.

    Reading an image means OCR — there is no other route to its text — so when
    OCR is unavailable this raises rather than returning empty. Returning []
    used to bubble up as the generic "appears to be empty or contains no
    extractable text", which is misleading: the file is fine, the server just
    cannot read images.
    """
    if not tesseract_available():
        reason = (
            "OCR is not available on this server"
            if not _tesseract_binary_present()
            else "OCR is disabled on this server"
        )
        raise UnreadableDocumentError(
            f"{filename} is an image, and {reason}, so no text can be read from it. "
            f"Please upload a text document instead ({TEXT_FORMATS_LABEL})."
        )

    text = ocr_image_bytes(content)
    stripped = text.strip()
    if len(stripped) < OCR_MIN_TEXT_LEN:
        raise UnreadableDocumentError(
            f"No readable text was found in {filename}. If it is a photo or "
            f"diagram, upload the source document instead."
        )
    # Prefix a section header that v2 chunker will detect (uppercase line)
    tagged = f"## OCR_EXTRACTED\n{stripped}"
    return [(1, tagged, True)]


def extract_pages_from_pptx(content: bytes) -> List[Tuple[int, str, bool]]:
    """One logical page per slide. Slide number becomes page_number."""
    if not _PPTX_AVAILABLE:
        raise RuntimeError("python-pptx is not installed (pip install python-pptx)")

    prs = _Presentation(io.BytesIO(content))
    pages: List[Tuple[int, str, bool]] = []
    for i, slide in enumerate(prs.slides, 1):
        title_text = ""
        body_lines: List[str] = []
        title_shape = getattr(slide.shapes, "title", None)
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            txt = (shape.text_frame.text or "").strip()
            if not txt:
                continue
            if title_shape is not None and shape == title_shape and not title_text:
                title_text = txt
            else:
                body_lines.append(txt)

        if not title_text and not body_lines:
            continue
        header = f"## Slide {i}: {title_text}" if title_text else f"## Slide {i}"
        page_text = header + ("\n" + "\n".join(body_lines) if body_lines else "")
        pages.append((i, page_text, False))
    return pages


# ---------------------------------------------------------------------------
# Unified entry point
# ---------------------------------------------------------------------------
def extract_pages_from_file(
    file_content: bytes,
    filename: str,
    stats: Optional[dict] = None,
) -> List[Tuple[int, str, bool]]:
    """
    Dispatch to the right extractor based on extension.

    PDF is handled by rag_service (it owns PyMuPDF + the OCR fallback logic)
    and is imported lazily here to avoid a circular import.

    ``stats``, when a caller passes a dict, is filled in with what the extractor
    learned about coverage (currently ``total_pages`` / ``unreadable_pages`` for
    PDFs) so the caller can tell the user that part of the document was skipped.
    Optional so existing callers are unaffected.
    """
    ext = get_file_type(filename)

    if ext == ".pdf":
        from app.services.rag_service import extract_pages_from_pdf_with_ocr
        return extract_pages_from_pdf_with_ocr(file_content, filename=filename, stats=stats)
    if ext == ".docx":
        return extract_pages_from_docx(file_content)
    if ext == ".xlsx":
        return extract_pages_from_xlsx(file_content, filename)
    if ext == ".xls":
        return extract_pages_from_xls(file_content, filename)
    if ext == ".csv":
        return extract_pages_from_csv(file_content, filename)
    if ext in _IMAGE_EXTS:
        return extract_pages_from_image(file_content, filename)
    if ext == ".pptx":
        return extract_pages_from_pptx(file_content)

    raise ValueError(f"Unsupported file type: {ext}")


# ---------------------------------------------------------------------------
# Legacy image-to-PDF conversion + media type helpers (unchanged)
# ---------------------------------------------------------------------------
async def convert_image_to_pdf(file_content: bytes, filename: str) -> tuple[bytes, str]:
    """Convert TIS/TIF/TIFF image to PDF (used by legacy upload path)."""
    try:
        image = Image.open(io.BytesIO(file_content))
        if image.mode not in ("RGB", "L", "1"):
            image = image.convert("RGB")
        img_bytes = io.BytesIO()
        image.save(img_bytes, format="PNG")
        img_bytes.seek(0)
        pdf_bytes = img2pdf.convert(img_bytes.read())
        new_filename = filename.rsplit(".", 1)[0] + ".pdf"
        return pdf_bytes, new_filename
    except Exception as conv_error:
        print(f"[ERROR] Error converting image to PDF: {conv_error}")
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Failed to convert image to PDF: {str(conv_error)}",
        )


def needs_image_conversion(filename: str) -> bool:
    if not filename or "." not in filename:
        return False
    return filename.lower().split(".")[-1] in ["tis", "tif", "tiff"]


def determine_media_type(filename: str) -> str:
    filename_lower = filename.lower()
    if filename_lower.endswith(".pdf"):
        return "application/pdf"
    if filename_lower.endswith(".png"):
        return "image/png"
    if filename_lower.endswith((".jpg", ".jpeg")):
        return "image/jpeg"
    if filename_lower.endswith(".gif"):
        return "image/gif"
    if filename_lower.endswith(".webp"):
        return "image/webp"
    if filename_lower.endswith((".tiff", ".tif")):
        return "image/tiff"
    if filename_lower.endswith((".doc", ".docx")):
        return "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    if filename_lower.endswith((".xls", ".xlsx")):
        return "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    if filename_lower.endswith(".csv"):
        return "text/csv"
    if filename_lower.endswith(".pptx"):
        return "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    if filename_lower.endswith(".txt"):
        return "text/plain"
    return "application/octet-stream"
