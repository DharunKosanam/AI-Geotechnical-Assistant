"""KB upload format-handler registry (Phase 2).

A registry keyed by file extension, mirroring the calculator registry pattern.
Each handler turns an uploaded file into a normalised list of
``(unit_number, text, ocr_flag)`` pages (the exact shape the v2 chunker already
consumes) PLUS format metadata, and carries its own extraction-validity check.

Scope: this registry is used ONLY by the KB upload feature (Phase 4). The live
thread/user upload path keeps using ``file_processing.extract_pages_from_file``
unchanged — nothing here touches it. PDF and DOCX reuse the existing extractors;
TXT, MD, PPTX (slides + speaker notes) and the metadata-indexed XLSX/CSV handlers
are added here.

Accepted for the KB: PDF, DOCX, TXT, MD, PPTX, XLSX, CSV. Images are REJECTED
(the KB is a text corpus; OCR'd images belong on the per-thread path), as is any
other format — both with a message listing what is accepted.

Spreadsheets are ROW-INDEXED (v3-xlsx): every non-empty row becomes one line
with its columns in order, sheet name + header repeated per block, capped at
XLSX_MAX_ROWS_PER_SHEET rows per sheet. The earlier metadata-only index
(headers + deduplicated labels) destroyed row associations — "who booked the
interrogator" was unanswerable even when retrieval found the right chunk.
"""
from __future__ import annotations

import io
import os
from dataclasses import dataclass, field
from typing import Any, Callable, Dict, List, Tuple

# Normalised extraction unit shared with the chunker: (unit_number, text, ocr).
Page = Tuple[int, str, bool]


class UnsupportedFormatError(ValueError):
    """Raised for a file the KB registry does not accept (images, unknown)."""


@dataclass
class HandlerResult:
    """What every handler returns.

    ``pages`` feed the chunker unchanged. ``format_metadata`` is per-format
    structural info (slide/sheet counts, headers, notes count) surfaced to the
    uploader and stored alongside provenance.
    """

    pages: List[Page]
    source_format: str
    format_metadata: Dict[str, Any] = field(default_factory=dict)

    def total_chars(self) -> int:
        return sum(len(t) for _, t, _ in self.pages)


@dataclass(frozen=True)
class FormatHandler:
    """One registered format. ``extract`` may raise on a hard failure (corrupt
    file / missing library); ``validate`` reports whether the *result* is usable
    (its own extraction-validity check), returning (ok, reason)."""

    ext: str          # includes the leading dot, e.g. ".pdf"
    label: str        # display label, e.g. "PDF"
    extract: Callable[[bytes, str], HandlerResult]
    validate: Callable[[HandlerResult], Tuple[bool, str]]


# ---------------------------------------------------------------------------
# Extension helpers
# ---------------------------------------------------------------------------
def _ext(filename: str) -> str:
    return os.path.splitext(filename or "")[1].lower()


# Images are explicitly recognised so we can reject them with a tailored message
# (they are accepted on the per-thread path, just not for the shared KB).
IMAGE_EXTENSIONS = frozenset({".png", ".jpg", ".jpeg", ".tif", ".tiff", ".gif", ".bmp", ".webp"})


# ---------------------------------------------------------------------------
# Validity checks (each handler's own extraction-validity gate)
# ---------------------------------------------------------------------------
def _validate_min_chars(min_chars: int) -> Callable[[HandlerResult], Tuple[bool, str]]:
    def _check(result: HandlerResult) -> Tuple[bool, str]:
        if not result.pages:
            return False, "no extractable content was found"
        n = result.total_chars()
        if n < min_chars:
            return False, f"only {n} characters extracted (need >= {min_chars})"
        return True, ""
    return _check


def _validate_pptx(result: HandlerResult) -> Tuple[bool, str]:
    slides = result.format_metadata.get("slides", 0)
    if slides == 0 or not result.pages:
        return False, "no slides with readable text were found"
    if result.total_chars() < 20:
        return False, f"only {result.total_chars()} characters across {slides} slide(s)"
    return True, ""


def _validate_spreadsheet(result: HandlerResult) -> Tuple[bool, str]:
    if not result.pages:
        return False, "no sheets or columns could be read"
    # Metadata-indexed: even a header-only sheet is valid (findable by structure).
    if result.total_chars() < 3:
        return False, "no headers or labels could be extracted"
    return True, ""


# ---------------------------------------------------------------------------
# Handlers
# ---------------------------------------------------------------------------
def _extract_pdf(content: bytes, filename: str) -> HandlerResult:
    # Reuse the canonical PyMuPDF + OCR-fallback extractor (lazy import avoids a
    # heavy module load and any import cycle).
    from app.services.file_processing import UnreadableDocumentError
    from app.services.rag_service import extract_pages_from_pdf_with_ocr

    try:
        triples: List[Page] = extract_pages_from_pdf_with_ocr(content, filename=filename)
    except UnreadableDocumentError:
        # The shared extractor now RAISES for a PDF with no readable text so the
        # chat upload chip can explain it. The KB pipeline classifies that same
        # condition itself, from an empty page list, as skip_reason="scanned" --
        # keep feeding it exactly that so KB behaviour is unchanged.
        triples = []
    ocr_pages = sum(1 for _, _, ocr in triples if ocr)
    return HandlerResult(
        pages=triples,
        source_format="pdf",
        format_metadata={"pages": len(triples), "ocr_pages": ocr_pages},
    )


def _extract_docx(content: bytes, filename: str) -> HandlerResult:
    from app.services.file_processing import extract_pages_from_docx

    triples: List[Page] = extract_pages_from_docx(content)
    return HandlerResult(
        pages=triples,
        source_format="docx",
        format_metadata={"pages": len(triples)},
    )


def _decode_text(content: bytes) -> str:
    """Best-effort decode: UTF-8 (with BOM), falling back to latin-1 so no byte
    sequence raises. Never lossy-throws."""
    for enc in ("utf-8-sig", "utf-8"):
        try:
            return content.decode(enc)
        except UnicodeDecodeError:
            continue
    return content.decode("latin-1", errors="replace")


def _extract_txt(content: bytes, filename: str) -> HandlerResult:
    text = _decode_text(content).strip()
    pages: List[Page] = [(1, text, False)] if text else []
    return HandlerResult(pages=pages, source_format="txt", format_metadata={"chars": len(text)})


def _extract_md(content: bytes, filename: str) -> HandlerResult:
    # Markdown is kept verbatim; its '#' headers help the v2 chunker's
    # section detection. One extraction unit; the chunker splits it.
    text = _decode_text(content).strip()
    pages: List[Page] = [(1, text, False)] if text else []
    return HandlerResult(pages=pages, source_format="md", format_metadata={"chars": len(text)})


def _extract_pptx(content: bytes, filename: str) -> HandlerResult:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as e:  # pragma: no cover - dependency guard
        raise RuntimeError("python-pptx is not installed (pip install python-pptx)") from e

    prs = Presentation(io.BytesIO(content))
    pages: List[Page] = []
    slides_with_notes = 0
    for i, slide in enumerate(prs.slides, 1):
        title_text = ""
        body_lines: List[str] = []
        title_shape = getattr(slide.shapes, "title", None)
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            txt = (shape.text_frame.text or "").strip()
            if not txt:
                continue
            if title_shape is not None and shape == title_shape and not title_text:
                title_text = txt
            else:
                body_lines.append(txt)

        # Speaker notes — the KB-specific addition over the live PPTX handler.
        notes = ""
        if getattr(slide, "has_notes_slide", False):
            ns = slide.notes_slide
            if ns is not None and ns.notes_text_frame is not None:
                notes = (ns.notes_text_frame.text or "").strip()
        if notes:
            slides_with_notes += 1

        if not title_text and not body_lines and not notes:
            continue
        header = f"## Slide {i}: {title_text}" if title_text else f"## Slide {i}"
        parts = [header]
        if body_lines:
            parts.append("\n".join(body_lines))
        if notes:
            parts.append(f"### Notes\n{notes}")
        pages.append((i, "\n".join(parts), False))

    return HandlerResult(
        pages=pages,
        source_format="pptx",
        format_metadata={"slides": len(pages), "slides_with_notes": slides_with_notes},
    )


def _extract_xlsx(content: bytes, filename: str) -> HandlerResult:
    # Row-structured extraction (v3-xlsx), shared with the live upload path.
    # Replaces the old Headers:/Labels: metadata index, which flattened every
    # sheet into one deduplicated cell list and made row associations ("who
    # booked what, when") unrecoverable at retrieval time.
    from app.services.file_processing import extract_xlsx_sheets

    pages, meta = extract_xlsx_sheets(content, filename)
    return HandlerResult(
        pages=pages,
        source_format="xlsx",
        format_metadata=meta,
    )


def _extract_csv(content: bytes, filename: str) -> HandlerResult:
    # Same row-structured path as .xlsx (see _extract_xlsx).
    from app.services.file_processing import extract_csv_rows

    pages, meta = extract_csv_rows(content, filename)
    return HandlerResult(
        pages=pages,
        source_format="csv",
        format_metadata=meta,
    )


# ---------------------------------------------------------------------------
# Registry
# ---------------------------------------------------------------------------
HANDLERS: Dict[str, FormatHandler] = {
    ".pdf": FormatHandler(".pdf", "PDF", _extract_pdf, _validate_min_chars(20)),
    ".docx": FormatHandler(".docx", "DOCX", _extract_docx, _validate_min_chars(20)),
    ".txt": FormatHandler(".txt", "TXT", _extract_txt, _validate_min_chars(20)),
    ".md": FormatHandler(".md", "MD", _extract_md, _validate_min_chars(20)),
    ".pptx": FormatHandler(".pptx", "PPTX", _extract_pptx, _validate_pptx),
    ".xlsx": FormatHandler(".xlsx", "XLSX", _extract_xlsx, _validate_spreadsheet),
    ".csv": FormatHandler(".csv", "CSV", _extract_csv, _validate_spreadsheet),
}

ACCEPTED_EXTENSIONS: Tuple[str, ...] = tuple(HANDLERS.keys())
ACCEPTED_LABEL = ", ".join(h.label for h in HANDLERS.values())


def get_handler(filename: str) -> "FormatHandler | None":
    """The handler for this file's extension, or None if not accepted."""
    return HANDLERS.get(_ext(filename))


def is_accepted(filename: str) -> bool:
    return _ext(filename) in HANDLERS


def extract_kb_document(content: bytes, filename: str) -> HandlerResult:
    """Extract via the registered handler, or raise UnsupportedFormatError with a
    message that lists what IS accepted. Images get a tailored rejection."""
    ext = _ext(filename)
    handler = HANDLERS.get(ext)
    if handler is None:
        if ext in IMAGE_EXTENSIONS:
            raise UnsupportedFormatError(
                f"Images are not accepted for the knowledge base. Accepted formats: {ACCEPTED_LABEL}."
            )
        shown = ext or (filename or "the file")
        raise UnsupportedFormatError(
            f"Unsupported file type '{shown}'. Accepted formats: {ACCEPTED_LABEL}."
        )
    return handler.extract(content, filename)
