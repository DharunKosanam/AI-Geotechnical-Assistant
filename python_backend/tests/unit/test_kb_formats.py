"""Phase 2: KB format-handler registry — one check per format + rejections."""
import io

import pytest

from app.services import kb_formats as kf


# --- fixtures: build tiny real files in-memory --------------------------------
def _docx_bytes():
    from docx import Document
    d = Document()
    d.add_heading("Consolidation Test Report", level=1)
    d.add_paragraph("This document describes the consolidation behaviour of soft clay.")
    buf = io.BytesIO()
    d.save(buf)
    return buf.getvalue()


def _pptx_bytes():
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Slope Stability"
    slide.placeholders[1].text = "Factor of safety exceeds 1.4 for the west face"
    slide.notes_slide.notes_text_frame.text = "Speaker note: remember the perched water table"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _xlsx_bytes():
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Boreholes"
    ws.append(["Depth (m)", "cu (kPa)", "Description"])
    ws.append([1.0, 96, "stiff clay"])
    ws.append([2.0, 110, "very stiff clay"])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _pdf_bytes():
    import fitz
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text(
        (72, 72),
        "Geotechnical site investigation report describing pile foundation design "
        "and consolidation settlement analysis for a soft clay deposit.",
        fontsize=11,
    )
    data = doc.tobytes()
    doc.close()
    return data


# --- text formats -------------------------------------------------------------
def test_txt():
    r = kf.extract_kb_document(b"Soil mechanics lecture notes on effective stress.", "notes.txt")
    assert r.source_format == "txt"
    assert "effective stress" in r.pages[0][1]
    assert kf.get_handler("notes.txt").validate(r)[0] is True


def test_md():
    r = kf.extract_kb_document(b"# Erosion\n\nInternal erosion mechanisms in embankment dams.", "e.md")
    assert r.source_format == "md"
    assert r.pages[0][1].startswith("# Erosion")  # markdown kept verbatim
    assert kf.get_handler("e.md").validate(r)[0] is True


# --- office formats -----------------------------------------------------------
def test_docx():
    r = kf.extract_kb_document(_docx_bytes(), "report.docx")
    assert r.source_format == "docx"
    assert "consolidation behaviour" in r.pages[0][1].lower()
    assert kf.get_handler("report.docx").validate(r)[0] is True


def test_pptx_includes_slides_and_notes():
    r = kf.extract_kb_document(_pptx_bytes(), "deck.pptx")
    text = "\n".join(p[1] for p in r.pages)
    assert "Slope Stability" in text            # title
    assert "Factor of safety" in text            # body
    assert "perched water table" in text         # SPEAKER NOTES (KB-specific)
    assert r.format_metadata["slides_with_notes"] >= 1
    assert kf.get_handler("deck.pptx").validate(r)[0] is True


def test_xlsx_is_metadata_indexed_not_cell_dump():
    r = kf.extract_kb_document(_xlsx_bytes(), "bh.xlsx")
    text = "\n".join(p[1] for p in r.pages)
    assert "Boreholes" in text                          # sheet name
    assert "Depth (m)" in text and "Description" in text  # headers
    assert "stiff clay" in text                          # text label
    assert "96" not in text and "110" not in text        # numeric cells NOT dumped
    assert r.format_metadata["sheets"] == ["Boreholes"]


def test_csv_is_metadata_indexed():
    r = kf.extract_kb_document(b"site,depth,soil\nAlpha,1.0,clay\nBeta,2.0,sand\n", "d.csv")
    text = r.pages[0][1]
    assert "Columns: site | depth | soil" in text
    assert "clay" in text and "sand" in text
    assert "1.0" not in text  # numeric cell not dumped
    assert r.format_metadata["columns"] == ["site", "depth", "soil"]
    assert r.format_metadata["sampled_rows"] == 2


def test_pdf():
    r = kf.extract_kb_document(_pdf_bytes(), "report.pdf")
    assert r.source_format == "pdf"
    joined = "\n".join(p[1] for p in r.pages).lower()
    assert "pile foundation" in joined
    assert kf.get_handler("report.pdf").validate(r)[0] is True


# --- rejections ---------------------------------------------------------------
def test_image_rejected_with_message():
    with pytest.raises(kf.UnsupportedFormatError) as e:
        kf.extract_kb_document(b"\x89PNG...", "figure.png")
    assert "Images are not accepted" in str(e.value)
    assert kf.ACCEPTED_LABEL in str(e.value)


def test_unsupported_rejected_with_accepted_list():
    with pytest.raises(kf.UnsupportedFormatError) as e:
        kf.extract_kb_document(b"PK...", "archive.zip")
    assert "Accepted formats" in str(e.value)
    assert "PDF" in str(e.value) and "PPTX" in str(e.value)


def test_xls_not_accepted_for_kb():
    # Legacy .xls is intentionally out of the KB accepted set (spec lists XLSX/CSV).
    assert kf.is_accepted("old.xls") is False


# --- registry -----------------------------------------------------------------
def test_registry_shape():
    assert set(kf.ACCEPTED_EXTENSIONS) == {".pdf", ".docx", ".txt", ".md", ".pptx", ".xlsx", ".csv"}
    assert kf.is_accepted("a.PDF") is True   # case-insensitive
    assert kf.get_handler("a.bogus") is None


def test_validity_rejects_too_short():
    r = kf.extract_kb_document(b"hi", "tiny.txt")
    ok, reason = kf.get_handler("tiny.txt").validate(r)
    assert ok is False and "characters" in reason
